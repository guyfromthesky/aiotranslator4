
import copy
import os

import multiprocessing
import queue 
#from multiprocessing import Process, Queue, Pool
from multiprocessing import Process , Queue, Manager
import subprocess
import time
from datetime import datetime



# Functions used for processor
def check_list(item, list):
	if list == [] or list == [""]:
		return True
	for text in list:
		if item == text:
			return True
	return False

# Functions used for processor
def generate_sheet_list(sheet_list, translate_list):
	"""
	Function returns a sheet index List.
	Translate classified sheets by name or sheet index (min-max)
	"""
	if translate_list == [] or translate_list == [""]:
		return sheet_list
	#TotalSheet = len(sheet_list)
	index_list = []
	for sheet in translate_list:
		print("Sheet: ", sheet)
		if sheet.isnumeric():
			index_list.append(int(sheet))
		else:
			temp_list = sheet.split('-')
			#print('temp_list: ', temp_list)
			#print('Len temp_list: ', len(temp_list))
			invalid = False
			if len(temp_list) == 2:
				for temp_index in temp_list:
					#print("Item: ",  temp_index)
					if temp_index.isnumeric() != True:
						invalid = True
			else:
				invalid = True
			# Decide which sheet index is the lowest and highest to select the index between
			# Ex: 2,4 -> Translate from sheet 2 to sheet 4
			if invalid == False:
				
				A = int(temp_list[0])
				B = int(temp_list[1])
				
				if A < B:
					Min = A
					Max = B
				else:
					Min = B
					Max = A	
				#print('Min max: ', Min, Max)
				i = Min
				while i >= Min and i <= Max:
					index_list.append(i)
					i += 1
			else:
				i = 0
				for source in sheet_list:
					if source == sheet:
						index_list.append(i)		
					i += 1

	to_return = []

	i = 1
	for to_translate in sheet_list:
		for num in index_list:
			if i == num:
				to_return.append(to_translate)
		i+=1
	
	return to_return

def show_progress(cp, total_process):
	#os.system('CLS') 
	percent = int(1000 * cp / total_process)
	#print("Current progress: " +  str(cp) + '/ ' + str(total_process))
	return percent

#Para Translate
def para_translate(tasks_to_accomplish, tasks_that_are_done, MyTranslator):
	while True:
		#print('para_translate is running')
		try:
			Task = tasks_to_accomplish.get_nowait()
		except queue.Empty:
			break
		else:
			if Task != None:
				text_to_translate = Task.Text
				Task.Fail = 0
				#print('Current text: ', text_to_translate)
				source_texts = text_to_translate.split('\n')
				translated_text = MyTranslator.translate(source_texts)
				#print('Para Translated', Translated)
				i = 0
				for string in translated_text:
					if string == False:
						translated_text[i] = source_texts[i]
						#print('Fail to translate....')
						Task.Fail += 1
						i += 1
			
				merged_text = "\n".join(translated_text)
				Task.Text = merged_text
				tasks_that_are_done.put(Task)
			else:
				tasks_that_are_done.put(None)
	return True

def cell_translate(MyTranslator, tasks_to_accomplish):
	List = []
	Length = []
	for task in tasks_to_accomplish:
		Input = task.Text
		
		if isinstance(Input, list):
			Length.append(len(Input))
			List += Input
		elif isinstance(Input, str):
			Length.append(1)
			List.append(Input)

	#print('Translate cell: ', List)
	Result = MyTranslator.translate(List)
	#print('Result translate cell: ', Result)
	
	#for i in range(len(List)):
	#	print(List[i], ":", Result[i])

	current_index = 0
	for i in range(len(tasks_to_accomplish)):
		temp_result = Result[current_index:current_index+Length[i]]
		tasks_to_accomplish[i].Text = temp_result
		current_index+=Length[i]
	
	#print('tasks_to_accomplish', tasks_to_accomplish)
	
	return tasks_to_accomplish



# xlsx object
# Add a cell as a task with an address
class CellData:
	def __init__(self, SheetName, CellAddress, Text):
		self.Sheet = SheetName
		self.Address = CellAddress
		self.Text = Text
		self.Fail = 0

# pptx object
# Add a Paragraph as a task with an address
class TextFrameData:
	def __init__(self, paragraphs_num, runs_num, string):
		self.Paragraphs = paragraphs_num
		self.Runs = runs_num
		self.Text = string

# docx object

class DocxData:
	def __init__(self, ParagraphData, TableData):
		self.Paragraphs = ParagraphData
		self.Tables = TableData
	

class ParagraphsData:
	def __init__(self, Paragraph_Index, Run_Index, Text):
		self.Paragraph = Paragraph_Index
		self.Run = Run_Index
		self.Text = Text

class TableData:
	def __init__(self, paragraphs_num, runs_num, string):
		self.Paragraphs = paragraphs_num
		self.Runs = runs_num
		self.Text = string

class RunData:
	def __init__(self, runs_num, string):
		self.Runs = runs_num
		self.Text = string

#Workbook handle
def translate_workbook(progress_queue=None, result_queue=None, status_queue=None, MyTranslator=None, Options = {}):
	""" How Document Translator tool proccesses the translation in Excel file """
	print(locals())
	from openpyxl import load_workbook, worksheet, Workbook
	from openpyxl.styles import Font

	if 'SourceDocument' not in Options:
		status_queue.put('No source document input')
		return False
	else:
		SourceDocument = Options['SourceDocument']

	if 'OutputDocument' not in Options:
		now = datetime.now()
		timestamp = str(int(datetime.timestamp(now)))
		output_dir = os.path.dirname(SourceDocument)
		base_name = os.path.basename(SourceDocument) # File name
		source_name = os.path.splitext(base_name)[0] # File extension
		OutputDocument = output_dir + '/' + 'Translated_' + source_name + '_' + timestamp + '.xlsx'
	else:
		OutputDocument = Options['OutputDocument']

	### TOOL OPTIONS START ###
	# Check if an option is selected
	if 'TranslateSheetName' not in Options:
		TranslateSheetName = True
	else:
		TranslateSheetName = Options['TranslateSheetName']

	if 'Sheet' not in Options:
		Sheet = []
	else:
		Sheet = Options['Sheet']
	### TOOL OPTIONS END ###

	### TRANSLATE OPTIONS START ###
	# Check if an option is selected
	if 'DataOnly' not in Options:
		DataOnly = True
	else:
		DataOnly = Options['DataOnly']

	if 'SheetRemovalMode' not in Options:
		SheetRemovalMode = False
	else:
		SheetRemovalMode = Options['SheetRemovalMode']
	print('Data only:', DataOnly)
	try:
		if DataOnly == True:
			print('Data only')
			xlsx = load_workbook(SourceDocument, data_only=True)
		else:
			print('Disable data only')
			xlsx = load_workbook(SourceDocument)
	except Exception as e:
		status_queue.put('Failed to load the document: ' + str(e))
		return e
	### TRANSLATE OPTIONS END ###

	status_queue.put('Estimating total task to do...')
	progress_queue.put(0)
	#tasks_to_accomplish = Queue()
	#tasks_that_are_done = Queue()
	#processes = []
	
	task_list = []

	sheet_list = []
	for sheet in xlsx:
		sheet_list.append(sheet.title)
	#status_queue.put('Sheet List' + str(sheet_list))

	current_task_count = 0

	# 
	translate_list = generate_sheet_list(sheet_list, Sheet)	
	#status_queue.put('Translate List' + str(translate_list))
	# Check if there's a value in the cell, number of task increases by 1
	for sheet in translate_list:
		ws = xlsx[sheet]
		for row in ws.iter_rows():
			for cell in row:
				if cell.value != None:
					current_task_count+=1
	
	total_task_count = current_task_count
	cp = 0
	if current_task_count == 0:
		print('Done')
	else:
		status_queue.put('Total task: ' + str(total_task_count))
		percent = show_progress(cp, total_task_count)
		progress_queue.put(percent)
		memory_translation = 0
		fail_request = 0
		empty_cell = 0
		status_queue.put('Checking task detail...')

		status_queue.put('Pre-processing document...')
		for sheet in xlsx:
			if check_list(sheet.title, translate_list):
				status_queue.put("Checking sheet: " + str(sheet.title))
				for row in sheet.iter_rows():
					for cell in row:
						if cell.value != None:
							current_string = str(cell.value)
							# ValidateSourceText in aiotranslator lib
							result = MyTranslator.ValidateSourceText(current_string)
							if result == False:
								fail_request+=1
							elif result != True:
								cell.value = result
								memory_translation+=1
							else:
								list_string = current_string.split('\n')
								sheet_name = sheet.title
								cell_address = cell.column_letter + str(cell.row)
								Task = CellData(sheet_name, cell_address, list_string)
								task_list.append(Task)
				cp = fail_request + memory_translation
				percent = show_progress(cp, total_task_count)
				progress_queue.put(percent)
			else:
				if SheetRemovalMode:
					std = xlsx.get_sheet_by_name(sheet.title)
					xlsx.remove_sheet(std)

		status_queue.put('Empty cell: ' + str(empty_cell))
		status_queue.put('Fail request: ' + str(fail_request))
		status_queue.put('Translated by Memory: ' + str(memory_translation))
		remaining_task_count = total_task_count - cp	
		status_queue.put('Remained task: ' + str(remaining_task_count))
		
		Task_todo = []
		
		while len(task_list) > 0:
			Translated = []
			TaskLength = 0
			
			while TaskLength < 2000:
				if len(task_list) > 0:
					Input = task_list[0].Text
					if isinstance(Input, list):
						TempLen = TaskLength
						for tempString in Input:
							TempLen += len(tempString)
			
					elif isinstance(Input, str):
						TempLen = TaskLength + len(task_list[0].Text)
						
					if TempLen < 2000:
						TaskLength = TempLen
						Task_todo.append(task_list[0])
						del task_list[0]
					else:
						break	
				else:
					break
		
			Translated = cell_translate(MyTranslator, Task_todo)
			
			for task in Translated:
				Return = task
				NewSheet = Return.Sheet
				NewAdd = Return.Address
				NewVal = Return.Text
				if isinstance(NewVal, list):
					NewVal = '\r\n'.join(NewVal)

				ws = xlsx[NewSheet]
				cell  = ws[NewAdd]
				cell.value = str(NewVal)
				tempFont = copy.copy(cell.font)  
				tempFont.name = 'Times New Roman'
				cell.font = tempFont
		
			remaining_task_count = len(task_list)
			Message = str(remaining_task_count) + ' tasks remain....'
			status_queue.put(Message)	

			cp = total_task_count - remaining_task_count
			percent = show_progress(cp, total_task_count)
			progress_queue.put(percent)
			del Task_todo
			Task_todo = []

		if TranslateSheetName:
			status_queue.put('Translating sheet name...')
			to_translate = []
			#print('Translate list: ', translate_list)
			for sheet in xlsx:
				if check_list(sheet.title, translate_list):
					CurrentSheetName = sheet.title
					to_translate.append(CurrentSheetName)
			translated = MyTranslator.translate(to_translate)
			#print('Translate translated list: ', translated)
			index = 0
			for sheet in xlsx:
				if check_list(sheet.title, translate_list):
					try:
						sheet.title = translated[index][0:29]
					except:
						pass	
					index+=1
		status_queue.put('Exporting result....')	
		print('Exporting file to ', OutputDocument)
		
		try:
			xlsx.save(OutputDocument)
			if os.path.isfile(OutputDocument):
				return True
			else:
				return False	
		except Exception as e:
			status_queue.put('Failed to save the result: ' + str(e))
			return e
	status_queue.put('No thing to do with this file.')	


#**********************************************************************************
#pptx handle **********************************************************************
#**********************************************************************************

def translate_presentation(progress_queue=None, result_queue=None, status_queue=None, MyTranslator=None, Options = {}):
	from pptx import Presentation
	
	
	progress_queue.put(0)
	
	if not 'SourceDocument' in Options:
		status_queue.put('No source document input')
		return False
	else:
		SourceDocument = Options['SourceDocument']
	#print('SourceDocument: ', SourceDocument)
	if not 'OutputDocument' in Options:
		now = datetime.now()
		timestamp = str(int(datetime.timestamp(now)))
		output_dir = os.path.dirname(SourceDocument)
		base_name = os.path.basename(SourceDocument)
		source_name, ext = os.path.splitext(base_name)
		OutputDocument = output_dir + '/' + 'Translated_' + source_name + '_' + timestamp + ext
	else:
		OutputDocument = Options['OutputDocument']

	if not 'Multiple' in Options:
		Multiple = 8
	else:
		Multiple = Options['Multiple']
	
	pptx = Presentation(SourceDocument)
	print('Estimating total task to do...')
	status_queue.put('Estimating total task to do...')


	if Multiple == None:
		print('Turbo mode is not available for the presentation translator')
	Start = time.time()
	total = 0
	
	for slide in pptx.slides:
		for shape in slide.shapes:
			if shape.has_text_frame:
				for paragraph in shape.text_frame.paragraphs:
					string = paragraph.text
					if string != "":
						if string != None:
							if isinstance(string, str) == True:
									total+= 1
			if shape.has_table:	
				for row in shape.table.rows:
					for cell in row.cells:
						for paragraph in cell.text_frame.paragraphs:
							string = paragraph.text
							if string != "":
								if string != None:
									if isinstance(string, str) == True:
										total+= 1
	total_task_count = total
	print('Total task: ', total_task_count)
	status_queue.put('Total task: ' + str(total_task_count))
	cp = 0
	
	for slide in pptx.slides: 
		for shape in slide.shapes:
			if shape.has_text_frame:
				for paragraph in shape.text_frame.paragraphs:
					FirstRun = None
					current_text = ""
					for j in range(len(paragraph.runs)):
						run = paragraph.runs[j]
						if FirstRun == None:
							FirstRun = j		
						if run.text == '':
							if current_text != '':
								list_text = current_text.split('\n')
								Translated = MyTranslator.translate(list_text)
								Translated = "\n".join(Translated)
								Translated = Translated.replace("\r\r\n", "\r\n")
								paragraph.runs[FirstRun].text = Translated
								
							FirstRun = None
							current_text = ""
						else:
							current_text += str(run.text)
							run.text = ""

					if current_text != '':
						list_text = current_text.split('\n')
						Translated = MyTranslator.translate(list_text)
						Translated = "\n".join(Translated)
						Translated = Translated.replace("\r\r\n", "\r\n")
						paragraph.runs[FirstRun].text = Translated
				
					cp+=1
					percent = show_progress(cp, total_task_count)
					progress_queue.put(percent)

			if shape.has_table:	
				for row in shape.table.rows:
					for cell in row.cells:
						for paragraph in cell.text_frame.paragraphs:
							FirstRun = None
							current_text = ""
							for j in range(len(paragraph.runs)):
								run = paragraph.runs[j]
								if FirstRun == None:
									FirstRun = j		
								if run.text == '':
									if current_text != '':
										list_text = current_text.split('\n')
										Translated = MyTranslator.translate(list_text)
										Translated = "\n".join(Translated)
										Translated = Translated.replace("\r\r\n", "\r\n")
										paragraph.runs[FirstRun].text = Translated
										
									FirstRun = None
									current_text = ""
								else:
									current_text += str(run.text)
									run.text = ""

							if current_text != '':
								list_text = current_text.split('\n')
								Translated = MyTranslator.translate(list_text)
								Translated = "\n".join(Translated)
								Translated = Translated.replace("\r\r\n", "\r\n")
								paragraph.runs[FirstRun].text = Translated

							cp+=1
							percent = show_progress(cp, total_task_count)
							progress_queue.put(percent)

	percent = show_progress(total_task_count, total_task_count)
	progress_queue.put(percent)
	status_queue.put('Exporting document...')
	try:
		pptx.save(OutputDocument)
		if (os.path.isfile(OutputDocument)):
			End = time.time()
			Message = "Total time spend: " + str(int(End-Start)) + ' seconds.'
			status_queue.put(Message)
			return True
		else:
			return False	
	except Exception as e:
		status_queue.put('Error message: ' + str(e))
		return e

def translateFrametable(cell, MyTranslator, CurrentTask, total_task_count):
	from docx.shared import RGBColor

	for paragraph in cell.text_frame.paragraphs:
		FirstRun = None
		current_text = ""
		for j in range(len(paragraph.runs)):
			run = paragraph.runs[j]
			if FirstRun == None:
				FirstRun = j		

			if run.text == '':
				if current_text != '':
					#print('current_text ', current_text)
					list_text = current_text.split('\n')
					#print('list_text ', list_text)
					Translated = MyTranslator.translate(list_text)
					Translated = "\n".join(Translated)
					#Translated = Translated.replace("\r\r\n", "\r\n")
					#print('Translated ', Translated)
					paragraph.runs[FirstRun].text = Translated
				FirstRun = None
				current_text = ""
			else:
				current_text += str(run.text)
				run.text = ""
		if current_text != '':
			#print('current_text ', current_text)
			list_text = current_text.split('\n')
			#print('list_text ', list_text)
			Translated = MyTranslator.translate(list_text)
			Translated = "\r\n".join(Translated)
			Translated = Translated.replace("\r\r\n", "\r\n")
			#print('Translated ', Translated)
			paragraph.runs[FirstRun].text = Translated

		CurrentTask+=1

	if cell.has_table:
		for row in cell.table.rows:
			for newcell in row.cells:
				newcell = translateFrametable(cell, MyTranslator, CurrentTask, total_task_count)
				CurrentTask+= 1
		print('Another table inside a table cell.')

	return [newcell, CurrentTask]

#**********************************************************************************
#Docx handle **********************************************************************
#**********************************************************************************

def check_paragraph_style(paragraph, Debug = False):
	if len(paragraph.runs) <= 1:
		return False
	for i in range(len(paragraph.runs)):
		
		run = paragraph.runs[i]
		rbold = str(run.bold)
		ritalic = str(run.italic)
		rcolor = str(run.font.color.rgb)
		runderline = str(run.underline)
		if Debug == True:
			#print("Checking: ", paragraph.text)
			print('bold', run.bold)
			print('italic', run.italic)
			print('color', run.font.color.rgb)
			print('underline', run.underline)

		if rbold != "None" and rbold != "False":
			return True
		elif ritalic != "None" and ritalic != "False":
			return True
		elif rcolor != "000000" and rcolor != "None":
			return True
		elif runderline != "None" and runderline != "False":
			return True

	return False

def translate_docx(progress_queue=None, result_queue=None, status_queue=None, MyTranslator=None, Options = {}):
	import docx
	
	from docx.shared import RGBColor
	print('Estimating total work to do...')
	status_queue.put('Estimating total task to do...')
	progress_queue.put(0)
	
	if not 'SourceDocument' in Options:
		status_queue.put('No source document input')
		return False
	else:
		SourceDocument = Options['SourceDocument']

	if not 'OutputDocument' in Options:
		now = datetime.now()
		timestamp = str(int(datetime.timestamp(now)))
		output_dir = os.path.dirname(SourceDocument)
		base_name = os.path.basename(SourceDocument)
		source_name, ext = os.path.splitext(base_name)
		OutputDocument = output_dir + '/' + 'Translated_' + source_name + '_' + timestamp + ext
	else:
		OutputDocument = Options['OutputDocument']

	Mydocx = docx.Document(SourceDocument)
	
	Start = time.time()
	
	cp = 0
	total_task_count = 0
	for i in range(len(Mydocx.paragraphs)):
		total_task_count += 1
	for table in Mydocx.tables:	
		total_task_count += estimate_table(table)
	status_queue.put('Total task: ' + str(total_task_count))

	paragraph_data_list = []

	for i in range(len(Mydocx.paragraphs)):
		paragraph = Mydocx.paragraphs[i]
		if paragraph.text != '':
			#print('TempText', paragraph.text)
			FirstRun = None
			current_text = ""
			for j in range(len(paragraph.runs)):
				run = paragraph.runs[j]
				if FirstRun == None:
					FirstRun = j		

				if run.text == '':
					if current_text != '':
						ParagraphData = ParagraphsData(i, FirstRun, current_text)
						paragraph_data_list.append(ParagraphData)
					FirstRun = None
					current_text = ""
				else:
					current_text += str(run.text)
					run.text = ""
			
			if current_text != '':
				ParagraphData = ParagraphsData(i, FirstRun, current_text)
				paragraph_data_list.append(ParagraphData)
			#('Adding:', ParagraphData.Text)
	

	RemainedTask = len(paragraph_data_list)
	Task_todo = []
	#print('RemainedTask', RemainedTask)
	#print('paragraph_data_list', paragraph_data_list)
	
	while len(paragraph_data_list) > 0:
		Translated = []
		TaskLength = 0
			
		while TaskLength < 2000:
			if len(paragraph_data_list) > 0:
				Input = paragraph_data_list[0].Text
				if isinstance(Input, list):
					TempLen = TaskLength
					for tempString in Input:
						TempLen += len(tempString)
		
				elif isinstance(Input, str):
					TempLen = TaskLength + len(paragraph_data_list[0].Text)
					
				if TempLen < 2000:
					TaskLength = TempLen
					#print("TempLen", TempLen)
					Task_todo.append(paragraph_data_list[0])
					del paragraph_data_list[0]
				else:
					break	
			else:
				break

		Translated = cell_translate(MyTranslator, Task_todo)

		for task in Translated:
			Return = task
			Current_Par = Return.Paragraph
			Current_Run = Return.Run
			Current_Text = Return.Text
			Mydocx.paragraphs[Current_Par].runs[Current_Run].text = Current_Text

		
		RemainedTask = len(paragraph_data_list)
		cp += RemainedTask
		percent = show_progress(cp, total_task_count)
		progress_queue.put(percent)
		del Task_todo
		Task_todo = []
	
	for table in Mydocx.tables:

		result = translate_table(table, MyTranslator, cp, total_task_count)

		table = result[0]
		cp = result[1]
		percent = show_progress(cp, total_task_count)
		progress_queue.put(percent)

	percent = show_progress(total_task_count, total_task_count)
	progress_queue.put(percent)
	status_queue.put('Exporting...')
	print('Exporting')
	try:
		Mydocx.save(OutputDocument)
		if (os.path.isfile(OutputDocument)):
			End = time.time()
			Message = "Total time spend: " + str(int(End-Start)) + ' seconds.'
			status_queue.put(Message)
			return True
		else:
			status_queue.put('Fail to export file')
			return False
	except Exception as e:
		status_queue.put('Error message: ' + str(e))
		return e

def estimate_table(intable):
	total = 0
	for row in intable.rows:
		for cell in row.cells:
			for paragraph in cell.paragraphs:

				for run in paragraph.runs:
					runtext = run.text
					if runtext != "":
						if not runtext.isdigit():
							total+= 1
			for table in cell.tables:
				total+= estimate_table(table)			
	return total

def translate_table(intable, MyTranslator, CurrentTask, total_task_count, TaskPool = 4):
	from docx.shared import RGBColor

	#tasks_to_accomplish = Queue()
	#tasks_that_are_done = Queue()
	
	for row in intable.rows:
		for cell in row.cells:
			for paragraph in cell.paragraphs:
				FirstRun = None
				current_text = ""
				for j in range(len(paragraph.runs)):
					run = paragraph.runs[j]
					if FirstRun == None:
						FirstRun = j		

					if run.text == '':
						if current_text != '':
							#print('current_text ', current_text)
							list_text = current_text.split('\n')
							#print('list_text ', list_text)
							Translated = MyTranslator.translate(list_text)
							Translated = "\n".join(Translated)
							#Translated = Translated.replace("\r\r\n", "\r\n")
							#print('Translated ', Translated)
							paragraph.runs[FirstRun].text = Translated
						FirstRun = None
						current_text = ""
					else:
						current_text += str(run.text)
						run.text = ""
				if current_text != '':
					#print('current_text ', current_text)
					list_text = current_text.split('\n')
					#print('list_text ', list_text)
					Translated = MyTranslator.translate(list_text)
					Translated = "\n".join(Translated)
					#Translated = Translated.replace("\r\r\n", "\r\n")
					#print('Translated ', Translated)
					paragraph.runs[FirstRun].text = Translated

				CurrentTask+=1
				
			for table in cell.tables:
				print('Another table inside a table celll.')
				translate_table(table, MyTranslator, CurrentTask, total_task_count)

	return [intable, CurrentTask]

#**********************************************************************************
# UI handle ***********************************************************************
#**********************************************************************************
def translate_msg(progress_queue=None, result_queue=None, status_queue=None, MyTranslator=None, Options = {}):
	from outlook_msg import Message
	from docx import Document
	from docx.shared import Inches

	print('Estimating total work to do...')
	status_queue.put('Estimating total task to do...')
	progress_queue.put(0)
	
	if not 'SourceDocument' in Options:
		status_queue.put('No source document input')
		return False
	else:
		SourceDocument = Options['SourceDocument']

	if not 'OutputDocument' in Options:
		now = datetime.now()
		timestamp = str(int(datetime.timestamp(now)))
		output_dir = os.path.dirname(SourceDocument)
		base_name = os.path.basename(SourceDocument)
		source_name, ext = os.path.splitext(base_name)
		OutputDocument = output_dir + '/' + 'Translated_' + source_name + '_' + timestamp + ext
	else:
		OutputDocument = Options['OutputDocument']

	Start = time.time()
	
	cp = 0
	total_task_count = 0

	document = Document()

	with open(SourceDocument) as msg_file:
		msg = Message(msg_file)

	contents = msg.body
	subject = msg.subject
	translated = MyTranslator.translate(subject)	
	document.add_heading(translated, 0)

	para = contents.split('\n')
	total_task_count += len(para)
	cp = 0
	for Par in para:
		if Par not in ["", ' ', '\r', '\n']:
			translated = MyTranslator.translate(Par)
			document.add_paragraph(translated)
			cp+=1
			show_progress(cp, total_task_count)
	try:
		document.save(OutputDocument + '.docx')
		if (os.path.isfile(OutputDocument+ '.docx')):
			End = time.time()
			Message = "Total time spend: " + str(int(End-Start)) + ' seconds.'
			status_queue.put(Message)
			return True
		else:
			status_queue.put('Fail to export file')
			return False
	except Exception as e:
		status_queue.put('Error message: ' + str(e))
		return e